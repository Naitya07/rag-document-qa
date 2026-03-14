"""Document parsing and text chunking.

Supports: PDF, DOCX, PPTX, TXT, CSV, Markdown, HTML, JSON,
          Images (PNG, JPG, WEBP) via OCR,
          Audio/Video (MP3, WAV, MP4, MOV, WEBM) via Whisper transcription.
"""

import csv
import io
import json
import os
import tempfile

import fitz  # PyMuPDF


# ── Text-based extractors ───────────────────────────────────

def extract_text_from_pdf(file_bytes: bytes) -> list[dict]:
    """Extract text from a PDF file."""
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append({"page": i + 1, "text": text.strip()})
    doc.close()
    return pages


def extract_text_from_docx(file_bytes: bytes) -> list[dict]:
    """Extract text from a DOCX file."""
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    sections = []
    for i in range(0, len(paragraphs), 10):
        section_text = "\n".join(paragraphs[i:i + 10])
        if section_text.strip():
            sections.append({"page": i // 10 + 1, "text": section_text.strip()})
    return sections


def extract_text_from_pptx(file_bytes: bytes) -> list[dict]:
    """Extract text from a PPTX file (one section per slide)."""
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slide_text = "\n".join(texts)
        if slide_text.strip():
            slides.append({"page": i + 1, "text": slide_text.strip()})
    return slides


def extract_text_from_txt(file_bytes: bytes) -> list[dict]:
    """Extract text from a plain text or markdown file."""
    text = file_bytes.decode("utf-8", errors="ignore")
    lines = text.split("\n")
    sections = []
    for i in range(0, len(lines), 30):
        section_text = "\n".join(lines[i:i + 30])
        if section_text.strip():
            sections.append({"page": i // 30 + 1, "text": section_text.strip()})
    return sections


def extract_text_from_csv(file_bytes: bytes) -> list[dict]:
    """Extract text from a CSV file."""
    text = file_bytes.decode("utf-8", errors="ignore")
    reader = csv.reader(io.StringIO(text))
    rows = [", ".join(row) for row in reader if any(cell.strip() for cell in row)]
    sections = []
    for i in range(0, len(rows), 20):
        section_text = "\n".join(rows[i:i + 20])
        if section_text.strip():
            sections.append({"page": i // 20 + 1, "text": section_text.strip()})
    return sections


def extract_text_from_html(file_bytes: bytes) -> list[dict]:
    """Extract text from an HTML file."""
    from bs4 import BeautifulSoup
    soup = BeautifulSoup(file_bytes, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()
    text = soup.get_text(separator="\n", strip=True)
    lines = text.split("\n")
    sections = []
    for i in range(0, len(lines), 30):
        section_text = "\n".join(lines[i:i + 30])
        if section_text.strip():
            sections.append({"page": i // 30 + 1, "text": section_text.strip()})
    return sections


def extract_text_from_json(file_bytes: bytes) -> list[dict]:
    """Extract text from a JSON file by flattening it."""
    text = file_bytes.decode("utf-8", errors="ignore")
    data = json.loads(text)

    def flatten(obj, prefix=""):
        lines = []
        if isinstance(obj, dict):
            for k, v in obj.items():
                lines.extend(flatten(v, f"{prefix}{k}: "))
        elif isinstance(obj, list):
            for i, item in enumerate(obj):
                lines.extend(flatten(item, f"{prefix}[{i}] "))
        else:
            lines.append(f"{prefix}{obj}")
        return lines

    flat_lines = flatten(data)
    sections = []
    for i in range(0, len(flat_lines), 20):
        section_text = "\n".join(flat_lines[i:i + 20])
        if section_text.strip():
            sections.append({"page": i // 20 + 1, "text": section_text.strip()})
    return sections


# ── Image OCR ────────────────────────────────────────────────

def extract_text_from_image(file_bytes: bytes) -> list[dict]:
    """Extract text from an image using EasyOCR."""
    import easyocr
    reader = easyocr.Reader(["en"], gpu=False, verbose=False)
    results = reader.readtext(file_bytes)
    text = "\n".join([r[1] for r in results])
    if not text.strip():
        return [{"page": 1, "text": "(No text detected in image)"}]
    return [{"page": 1, "text": text.strip()}]


# ── Audio / Video transcription ──────────────────────────────

def extract_text_from_media(file_bytes: bytes, filename: str) -> list[dict]:
    """Transcribe audio/video using OpenAI Whisper (local)."""
    import whisper

    ext = filename.rsplit(".", 1)[-1].lower()
    with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as f:
        f.write(file_bytes)
        tmp_path = f.name

    try:
        ffmpeg_path = os.path.join(os.path.expanduser("~"), "bin", "ffmpeg")
        if os.path.exists(ffmpeg_path):
            os.environ["PATH"] = os.path.dirname(ffmpeg_path) + ":" + os.environ.get("PATH", "")

        model = whisper.load_model("base")
        result = model.transcribe(tmp_path)

        segments = result.get("segments", [])
        if not segments:
            return [{"page": 1, "text": result.get("text", "(No speech detected)")}]

        # Group segments into ~60-second sections
        sections = []
        current_text = []
        section_start = 0
        section_num = 1

        for seg in segments:
            current_text.append(seg["text"].strip())
            if seg["end"] - section_start >= 60:
                sections.append({
                    "page": section_num,
                    "text": " ".join(current_text),
                })
                current_text = []
                section_start = seg["end"]
                section_num += 1

        if current_text:
            sections.append({
                "page": section_num,
                "text": " ".join(current_text),
            })

        return sections
    finally:
        os.unlink(tmp_path)


# ── Router ───────────────────────────────────────────────────

EXTRACTORS = {
    # Text-based
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "pptx": extract_text_from_pptx,
    "txt": extract_text_from_txt,
    "md": extract_text_from_txt,
    "csv": extract_text_from_csv,
    "html": extract_text_from_html,
    "htm": extract_text_from_html,
    "json": extract_text_from_json,
    # Images (OCR)
    "png": extract_text_from_image,
    "jpg": extract_text_from_image,
    "jpeg": extract_text_from_image,
    "webp": extract_text_from_image,
    # Audio / Video (Whisper)
    "mp3": None,  # handled by extract_text_from_media
    "wav": None,
    "mp4": None,
    "mov": None,
    "webm": None,
    "m4a": None,
}

MEDIA_EXTENSIONS = {"mp3", "wav", "mp4", "mov", "webm", "m4a"}


def extract_text(file_bytes: bytes, filename: str) -> list[dict]:
    """Route to the correct extractor based on file extension."""
    ext = filename.rsplit(".", 1)[-1].lower()

    if ext in MEDIA_EXTENSIONS:
        return extract_text_from_media(file_bytes, filename)

    extractor = EXTRACTORS.get(ext)
    if not extractor:
        raise ValueError(f"Unsupported file type: .{ext}")
    return extractor(file_bytes)


def chunk_pages(pages: list[dict], chunk_size: int = 500, overlap: int = 50) -> list[dict]:
    """Split page/section texts into overlapping chunks for embedding."""
    chunks = []
    for page_info in pages:
        words = page_info["text"].split()
        start = 0
        idx = 0
        while start < len(words):
            end = start + chunk_size
            chunk_text = " ".join(words[start:end])
            chunks.append({
                "text": chunk_text,
                "page": page_info["page"],
                "chunk_index": idx,
            })
            start += chunk_size - overlap
            idx += 1
    return chunks
